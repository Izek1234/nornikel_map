"""Document parsing (PDF/DOCX/TXT/MD) and chunking."""

import hashlib
import io
import re

CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200
MAX_CHUNKS = 60  # cost guard for MVP


def parse_file(filename: str, content: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return _parse_pdf(content)
    if name.endswith(".docx"):
        return _parse_docx(content)
    if name.endswith(".pptx"):
        return _parse_pptx(filename, content)
    if name.endswith((".html", ".htm")):
        return _parse_html(content)
    if name.endswith((".txt", ".md", ".markdown", ".csv")):
        return _decode_text(content)
    raise ValueError(f"Неподдерживаемый формат: {filename}. Поддерживаются PDF, DOCX, PPTX, HTML, TXT, MD, CSV.")


def _parse_pdf(content: bytes) -> str:
    # Quick magic-byte check — real PDFs start with %PDF
    if not content.lstrip()[:4].startswith(b"%PDF"):
        # Maybe it's HTML or plain text disguised as PDF (e.g. from Yandex Disk)
        try:
            return _parse_html(content) if b"<html" in content[:512].lower() else _decode_text(content)
        except Exception:
            raise ValueError("Файл имеет расширение .pdf, но не является PDF-документом")

    try:
        import fitz  # PyMuPDF
    except ModuleNotFoundError:
        raise ModuleNotFoundError(
            "PyMuPDF не установлен. Выполните: pip install pymupdf"
        )

    try:
        doc = fitz.open(stream=content, filetype="pdf")
    except Exception as e:
        raise ValueError(f"Не удалось открыть PDF: {e}")

    pages = []
    for page in doc:
        blocks = page.get_text("blocks")
        # Sort blocks by vertical position, then horizontal to maintain reading order
        blocks.sort(key=lambda b: (b[1], b[0]))
        text = "\n".join(b[4].strip() for b in blocks if b[4].strip())
        if text.strip():
            pages.append(text)
    return "\n\n".join(pages)


def _parse_docx(content: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(content))
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(filename: str, content: bytes) -> str:
    import os
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(content))
    slides_list = []
    
    # Establish public directory for extraction
    public_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "frontend", "public", "extracted_images", filename))
    os.makedirs(public_dir, exist_ok=True)

    for s_idx, slide in enumerate(prs.slides):
        slide_parts = []
        img_idx = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text_frame.text.strip().replace("\n", " ") for c in row.cells if c.text_frame.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
            
            # Picture extraction
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    img_name = f"slide_{s_idx+1}_img_{img_idx+1}.{ext}"
                    img_path = os.path.join(public_dir, img_name)
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    # Insert visual image tag marker into slide text
                    slide_parts.append(f"[IMAGE: /extracted_images/{filename}/{img_name}]")
                    img_idx += 1
                except Exception as e:
                    print("Error extracting image:", e)

        if slide_parts:
            slides_list.append("\n".join(slide_parts))
    return "\n\n=== SLIDE ===\n\n".join(slides_list)


def _parse_html(content: bytes) -> str:
    from bs4 import BeautifulSoup
    html_text = _decode_text(content)
    soup = BeautifulSoup(html_text, "html.parser")
    # Remove scripts and styles
    for script in soup(["script", "style", "header", "footer", "nav"]):
        script.extract()
    return soup.get_text(separator="\n", strip=True)


def _decode_text(content: bytes) -> str:
    for enc in ("utf-8", "cp1251", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def clean_text(text: str) -> str:
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str) -> list[dict]:
    """Split into ~CHUNK_SIZE char chunks on paragraph/sentence boundaries with overlap."""
    text = clean_text(text)
    if not text:
        return []

    chunks = []
    start = 0
    n = len(text)
    while start < n and len(chunks) < MAX_CHUNKS:
        end = min(start + CHUNK_SIZE, n)
        if end < n:
            # prefer to break on paragraph, then sentence, then space
            window = text[start:end]
            brk = window.rfind("\n\n")
            if brk < CHUNK_SIZE // 2:
                brk = max(window.rfind(". "), window.rfind(".\n"))
            if brk < CHUNK_SIZE // 2:
                brk = window.rfind(" ")
            if brk > CHUNK_SIZE // 2:
                end = start + brk + 1
        piece = text[start:end].strip()
        if piece:
            chunks.append({
                "text": piece,
                "hash": hashlib.sha256(piece.encode("utf-8")).hexdigest(),
                "index": len(chunks),
            })
        if end >= n:
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return chunks
